"""Content inserter — fills the empty PPTX with generated content and images."""

import os
import re

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

from src.json_parser import parse_llm_json


def _get_theme_colors(design_system: dict, theme_name: str) -> dict:
    """Extract text and background colors from the design system."""
    theme = design_system["themes"][theme_name]
    scheme = design_system["color_schemes"][theme["color_scheme"]]
    return {
        "text": RGBColor.from_string(scheme["text"].lstrip("#")),
        "background": RGBColor.from_string(scheme["background"].lstrip("#")),
    }


# ──────────────────────────────────────────────
# Value cleaning — flatten nested dicts/lists
# ──────────────────────────────────────────────


def _flatten_value(value) -> str:
    """
    Convert any LLM output value into clean plain text.
    Handles: str, list, dict, nested structures.
    """
    if isinstance(value, str):
        return _clean_text(value)

    if isinstance(value, list):
        # List of strings → bullet points
        cleaned = []
        for item in value:
            if isinstance(item, str):
                cleaned.append(f"• {_clean_text(item)}")
            elif isinstance(item, dict):
                # Nested dict inside a list
                for k, v in item.items():
                    cleaned.append(f"• {_clean_text(str(v))}")
            else:
                cleaned.append(f"• {_clean_text(str(item))}")
        return "\n".join(cleaned)

    if isinstance(value, dict):
        # Dict with sub-keys (e.g., {"left": [...], "right": [...]})
        parts = []
        for sub_key, sub_val in value.items():
            if isinstance(sub_val, list):
                for item in sub_val:
                    parts.append(f"• {_clean_text(str(item))}")
            elif isinstance(sub_val, str):
                parts.append(_clean_text(sub_val))
            else:
                parts.append(_clean_text(str(sub_val)))
        return "\n".join(parts)

    return _clean_text(str(value))


def _clean_text(text: str) -> str:
    """
    Strip JSON/Python artifacts and LLM junk from text.
    """
    text = text.strip()

    # Remove leading/trailing quotes that LLM sometimes adds
    if (text.startswith('"') and text.endswith('"')) or \
       (text.startswith("'") and text.endswith("'")):
        text = text[1:-1]

    # Remove markdown bold/italic markers
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)

    # Remove markdown bullet markers (we add our own)
    text = re.sub(r"^[-•]\s+", "", text)

    # Remove trailing backslash-n as literal text
    text = text.replace("\\n", "\n")

    # Strip <think>...</think> tags from LLM reasoning
    text = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL)

    # Remove JSON-like artifacts at start/end
    text = re.sub(r"^\{.*?['\"]:\s*['\"]", "", text)
    text = re.sub(r"['\"]?\}$", "", text)

    return text.strip()


# ──────────────────────────────────────────────
# Font size calculation
# ──────────────────────────────────────────────

_AVG_CHAR_WIDTH_FACTOR = 0.55
_LINE_HEIGHT_FACTOR = 1.2


def _estimate_font_size(
    text: str,
    box_width_emu: int,
    box_height_emu: int,
    max_font_pt: int,
    min_font_pt: int = 10,
) -> int:
    """
    Estimate the largest font size (in pt) that fits text inside a box.
    Steps down from max_font_pt until the text fits, or hits min_font_pt.
    """
    box_width_pt = box_width_emu / 12700
    box_height_pt = box_height_emu / 12700

    for font_pt in range(max_font_pt, min_font_pt - 1, -1):
        char_width = font_pt * _AVG_CHAR_WIDTH_FACTOR
        if char_width == 0:
            continue
        chars_per_line = max(1, int(box_width_pt / char_width))

        wrapped_lines = 0
        for segment in text.split("\n"):
            segment_len = len(segment)
            wrapped_lines += max(1, -(-segment_len // chars_per_line))

        line_height = font_pt * _LINE_HEIGHT_FACTOR
        available_lines = max(1, int(box_height_pt / line_height))

        if wrapped_lines <= available_lines:
            return font_pt

    return min_font_pt


def _get_max_font_pt(slot_key: str) -> int:
    """Determine the maximum (ideal) font size based on slot type."""
    key = slot_key.lower()
    if key in ("main_title", "title") or (key.startswith("title") and "sub" not in key):
        return 30
    if "subtitle" in key:
        return 24
    if "title" in key:
        return 22
    return 18


def insert_content_into_ppt(
    ppt_path: str,
    completed_slides: list,
    downloaded_images: list,
    image_requests: list,
    theme_name: str,
    design_system: dict,
    output_path: str = "final_presentation.pptx",
) -> str:
    """
    Insert generated content and images into the empty PPTX.

    Uses two layers of overflow protection:
    1. Proactive font scaling
    2. TEXT_TO_FIT_SHAPE as safety net
    """
    prs = Presentation(ppt_path)

    colors = _get_theme_colors(design_system, theme_name)
    text_color = colors["text"]
    bg_color = colors["background"]

    # Build slide_id → content map
    content_map = {}
    for slide in completed_slides:
        raw = slide["content"]
        content_map[slide["slide_id"]] = (
            parse_llm_json(raw) if isinstance(raw, str) else raw
        )

    # Build image lookup: (slide_id, slot_key) → image_path
    image_map = {}
    for img in downloaded_images:
        image_map[(img["slide_id"], img["slot_key"])] = img["image_path"]

    # Build set of ALL image slot keys
    image_slot_keys = set()
    for img in image_requests:
        image_slot_keys.add((img["slide_id"], img["slot_key"]))

    # Iterate slides
    for slide_index, slide in enumerate(prs.slides, start=1):
        if slide_index not in content_map:
            continue

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        slide_content = content_map[slide_index]

        for shape in list(slide.shapes):
            placeholder_text = ""

            if shape.has_text_frame:
                placeholder_text = shape.text_frame.text.strip()

            if not (placeholder_text.startswith("[") and placeholder_text.endswith("]")):
                continue

            slot_key = placeholder_text[1:-1]
            slot_id = (slide_index, slot_key)

            # ── IMAGE INSERTION ──
            if slot_id in image_slot_keys:
                if slot_id in image_map and os.path.exists(image_map[slot_id]):
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    slide.shapes.add_picture(
                        image_map[slot_id], left, top, width=width, height=height
                    )

                slide.shapes._spTree.remove(shape._element)
                continue

            # ── TEXT INSERTION ──
            if not shape.has_text_frame:
                continue

            if slot_key not in slide_content:
                continue

            tf = shape.text_frame
            raw_value = slide_content[slot_key]

            # ── CLEAN THE VALUE ──
            # Flatten nested dicts/lists into plain text
            value = _flatten_value(raw_value)

            if not value.strip():
                continue

            # Calculate the right font size BEFORE inserting
            max_font = _get_max_font_pt(slot_key)
            fitted_font_pt = _estimate_font_size(
                text=value,
                box_width_emu=shape.width,
                box_height_emu=shape.height,
                max_font_pt=max_font,
                min_font_pt=10,
            )
            font_size = Pt(fitted_font_pt)

            # Clear and configure text frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            if "\n" in value:
                lines = value.split("\n")
                p = tf.paragraphs[0]
                p.text = lines[0]
                p.font.size = font_size
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.JUSTIFY

                for line in lines[1:]:
                    bp = tf.add_paragraph()
                    bp.text = line
                    bp.font.size = font_size
                    bp.font.color.rgb = text_color
                    bp.alignment = PP_ALIGN.JUSTIFY
            else:
                p = tf.paragraphs[0]
                p.text = value
                p.font.size = font_size
                p.font.color.rgb = text_color
                p.alignment = PP_ALIGN.JUSTIFY

            # Remove dev-mode borders
            if shape.line:
                shape.line.fill.background()
            if shape.fill:
                shape.fill.background()

    prs.save(output_path)
    return output_path
