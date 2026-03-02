"""PPT maker — creates an empty PowerPoint with layout slots placed."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from src.models import PresentationPlan


def create_empty_ppt(
    planner_output: PresentationPlan,
    design_system: dict,
    output_path: str = "empty_template.pptx",
) -> str:
    """
    Create an empty PPTX with shapes and textboxes placed per layout.

    Parameters
    ----------
    planner_output : PresentationPlan
        The planner's output with theme and slide definitions.
    design_system : dict
        The loaded design system JSON.
    output_path : str
        Path to save the empty PPTX.

    Returns
    -------
    str
        Path to the saved PPTX file.
    """
    prs = Presentation()

    theme_name = planner_output.theme
    theme = design_system["themes"][theme_name]

    color_scheme = design_system["color_schemes"][theme["color_scheme"]]
    typography = design_system["typography"][theme["typography"]]

    for slide_plan in planner_output.slides:
        layout_name = slide_plan.layout
        layout_def = design_system["layouts"][layout_name]["slots"]

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

        for slot_key, slot in layout_def.items():
            slot_type = slot["type"]

            x = Inches(slot["x"])
            y = Inches(slot["y"])
            w = Inches(slot["w"])
            h = Inches(slot["h"])

            # ---------- SHAPES ----------
            if slot_type == "shape":
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(
                    color_scheme["primary"].replace("#", "")
                )
                shape.line.fill.background()
                continue

            # ---------- TEXT PLACEHOLDERS ----------
            textbox = slide.shapes.add_textbox(x, y, w, h)

            # Dev-mode visuals (light border + background)
            textbox.line.color.rgb = RGBColor(200, 200, 200)
            textbox.line.width = Pt(0.75)
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(245, 245, 245)
            textbox.fill.transparency = 0.85

            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            p = tf.paragraphs[0]

            # Placeholder label
            p.text = f"[{slot_key}]"

            # Base font
            p.font.name = typography["body_font"]
            p.font.size = Pt(typography["body_size"])
            p.font.color.rgb = RGBColor.from_string(
                color_scheme["text"].replace("#", "")
            )

            # Alignment
            align = slot.get("align", "left")
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            # Title override
            if slot_key == "title":
                p.font.name = typography["title_font"]
                p.font.size = Pt(typography["title_size"])

    prs.save(output_path)
    return output_path
